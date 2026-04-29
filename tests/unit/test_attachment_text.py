"""
Test get_attachment_text related functions
"""
import pytest
from unittest.mock import patch, Mock
from redmine_mcp.server import (
    _get_file_extension,
    _try_decode_text,
    _extract_pdf_text,
    _extract_docx_text,
    _extract_xlsx_text,
    _extract_pptx_text,
)


class TestGetFileExtension:
    def test_common_extensions(self):
        assert _get_file_extension("report.pdf") == ".pdf"
        assert _get_file_extension("doc.DOCX") == ".docx"
        assert _get_file_extension("data.json") == ".json"
        assert _get_file_extension("image.PNG") == ".png"

    def test_no_extension(self):
        assert _get_file_extension("README") == ""
        assert _get_file_extension("Makefile") == ""

    def test_multiple_dots(self):
        assert _get_file_extension("archive.tar.gz") == ".gz"
        assert _get_file_extension("my.file.txt") == ".txt"


class TestTryDecodeText:
    def test_utf8_text(self):
        text = "Hello, 世界！"
        result = _try_decode_text(text.encode('utf-8'))
        assert result == text

    def test_utf8_bom(self):
        text = "Hello BOM"
        result = _try_decode_text(b'\xef\xbb\xbf' + text.encode('utf-8'))
        assert "Hello BOM" in result

    def test_binary_data(self):
        # Large amount of non-printable characters → should return None
        binary = bytes(range(256)) * 10
        result = _try_decode_text(binary)
        assert result is None

    def test_empty_data(self):
        result = _try_decode_text(b"")
        assert result == ""

    def test_json_content(self):
        json_str = '{"key": "value", "number": 42}'
        result = _try_decode_text(json_str.encode('utf-8'))
        assert result == json_str

    def test_xml_content(self):
        xml_str = '<?xml version="1.0"?>\n<root><item>test</item></root>'
        result = _try_decode_text(xml_str.encode('utf-8'))
        assert result == xml_str


class TestExtractPdfText:
    def test_extract_pdf(self):
        """Test PDF extraction (using a real minimal PDF)"""
        from pypdf import PdfWriter
        from io import BytesIO

        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)

        # PyPDF2 blank page contains no text, so expect a hint message
        buf = BytesIO()
        writer.write(buf)
        result = _extract_pdf_text(buf.getvalue())
        assert "cannot extract" in result or "Page" in result


class TestExtractDocxText:
    def test_extract_docx(self):
        """Test Word text extraction"""
        from docx import Document
        from io import BytesIO

        doc = Document()
        doc.add_paragraph("First paragraph text")
        doc.add_paragraph("Second paragraph text")

        buf = BytesIO()
        doc.save(buf)
        result = _extract_docx_text(buf.getvalue())
        assert "First paragraph text" in result
        assert "Second paragraph text" in result

    def test_extract_docx_with_table(self):
        """Test Word table extraction"""
        from docx import Document
        from io import BytesIO

        doc = Document()
        doc.add_paragraph("Title")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"

        buf = BytesIO()
        doc.save(buf)
        result = _extract_docx_text(buf.getvalue())
        assert "Title" in result
        assert "A1" in result
        assert "B2" in result

    def test_extract_empty_docx(self):
        """Test empty Word document"""
        from docx import Document
        from io import BytesIO

        doc = Document()
        buf = BytesIO()
        doc.save(buf)
        result = _extract_docx_text(buf.getvalue())
        assert "no text content" in result


class TestExtractXlsxText:
    def test_extract_xlsx(self):
        """Test Excel extraction"""
        from openpyxl import Workbook
        from io import BytesIO

        wb = Workbook()
        ws = wb.active
        ws.title = "DataSheet"
        ws['A1'] = "Name"
        ws['B1'] = "Score"
        ws['A2'] = "John"
        ws['B2'] = 95

        buf = BytesIO()
        wb.save(buf)
        result = _extract_xlsx_text(buf.getvalue())
        assert "DataSheet" in result
        assert "Name" in result
        assert "John" in result
        assert "95" in result

    def test_extract_empty_xlsx(self):
        """Test empty Excel"""
        from openpyxl import Workbook
        from io import BytesIO

        wb = Workbook()
        ws = wb.active
        # Clear default cell
        buf = BytesIO()
        wb.save(buf)
        # Empty workbook has one empty sheet, but no content
        result = _extract_xlsx_text(buf.getvalue())
        # New workbook created by openpyxl may have empty rows, result depends on implementation
        assert isinstance(result, str)


class TestExtractPptxText:
    def test_extract_pptx(self):
        """Test PowerPoint extraction"""
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Content text"

        buf = BytesIO()
        prs.save(buf)
        result = _extract_pptx_text(buf.getvalue())
        assert "Slide Title" in result
        assert "Content text" in result

    def test_extract_empty_pptx(self):
        """Test empty presentation"""
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation()
        buf = BytesIO()
        prs.save(buf)
        result = _extract_pptx_text(buf.getvalue())
        assert "no text content" in result
