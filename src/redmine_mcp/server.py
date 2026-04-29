"""
Redmine MCP server main program
Provide MCP tools integrated with Redmine system
"""

import os
import argparse
from typing import Any
from datetime import datetime

# Ensure configuration is loaded before FastMCP initialization
# This will handle all environment variable settings, including FASTMCP_LOG_LEVEL
from .config import get_config
config = get_config()

from mcp.server.fastmcp import FastMCP
from mcp.server.fastmcp.utilities.types import Image
from .redmine_client import get_client, RedmineAPIError

# Create FastMCP server instance
mcp = FastMCP("Redmine MCP")


@mcp.tool()
def server_info() -> str:
    """Get server information and status"""
    config = get_config()
    return f"""Redmine MCP server started
- Redmine domain: {config.redmine_domain}
- Debug mode: {config.debug_mode}
- API timeout: {config.redmine_timeout}Second"""


@mcp.tool()
def health_check() -> str:
    """Health check tool to confirm that the server is functioning properly"""
    try:
        config = get_config()
        client = get_client()
        # Test connection
        if client.test_connection():
            return f"✓ The server is functioning normally and connected to {config.redmine_domain}"
        else:
            return f"✗ Unable to connect to Redmine server: {config.redmine_domain}"
    except Exception as e:
        return f"✗ Server exception: {str(e)}"


@mcp.tool()
def get_issue(issue_id: int, include_details: bool = True) -> str:
    """
    Get specific Redmine issue details
    
    Args:
        issue_id: Issue ID
        include_details: whether to include detailed information (description, notes, attachments, etc.)
    
    Returns:
        Detailed information about the issue, presented in an easy-to-read format
    """
    try:
        client = get_client()
        include_params = []
        if include_details:
            include_params = ['attachments', 'changesets', 'children', 'journals', 'relations', 'watchers']
        
        # Use the new get_issue_raw method to get complete information
        issue_data = client.get_issue_raw(issue_id, include=include_params)
        
        # Emoji correspondence table
        tracker_emojis = {'Bug': '🐛', 'defect': '🐛', 'Function': '✨', 'Feature': '✨', 'Refactor': '🔧', 'Refactor': '🔧', 'document': '📝', 'Documentation': '📝'}
        priority_emojis = {'Low': '🟢', 'normal': '🟡', 'high': '🟠', 'urgent': '🔴', 'immediate': '🔴'}
        status_emojis = {'Newly created': '📝', 'In practice': '🟡', 'responded': '🟡', 'Resolved': '🟢', 'Pending review': '🟢', 'ended': '⚫', 'Rejected': '⚫'}

        # Format basic information
        tracker_name = issue_data['tracker'].get('name', 'N/A')
        tracker_emoji = tracker_emojis.get(tracker_name, '📋')
        status_name = issue_data['status'].get('name', 'N/A')
        status_emoji = status_emojis.get(status_name, '🔵')
        priority_name = issue_data['priority'].get('name', 'N/A')
        priority_emoji = priority_emojis.get(priority_name, '🟡')

        assigned_name = issue_data.get('assigned_to', {}).get('name') if issue_data.get('assigned_to') else None
        assigned_to = f"@{assigned_name}" if assigned_name else 'Not assigned'
        category_info = issue_data.get('category', {}).get('name', 'Uncategorized') if issue_data.get('category') else 'Uncategorized'
        estimated = issue_data.get('estimated_hours')
        estimated_text = f"{estimated} hours" if estimated is not None else 'Not set'
        spent = issue_data.get('total_spent_hours') or issue_data.get('spent_hours')
        spent_text = f"{spent} hours" if spent is not None else 'Not recorded'

        # Extract specific custom fields
        custom_fields = issue_data.get('custom_fields', [])
        cf_map = {cf.get('id'): cf.get('value', '') for cf in custom_fields}
        actual_end_date = cf_map.get(23, 'Not set')
        resolve_date = cf_map.get(64, 'Not set')

        result = f"""Issue #{issue_data['id']}: {issue_data['subject']}

| Field | Value |
|------|------|
| **Issue ID** | #{issue_data['id']} |
| **Subject** | {issue_data['subject']} |
| **Project** | {issue_data['project'].get('name', 'N/A')} (ID: {issue_data['project'].get('id', 'N/A')}) |
| **Tracker** | `{tracker_emoji} {tracker_name}` |
| **Status** | `{status_emoji} {status_name}` |
| **Priority** | `{priority_emoji} {priority_name}` |
| **Category** | {category_info} |
| **Author** | {issue_data['author'].get('name', 'N/A')} |
| **Assigned to** | {assigned_to} |
| **Done ratio** | {issue_data.get('done_ratio', 0)}% |
| **Created on** | {issue_data.get('created_on', 'N/A')} |
| **Start date** | {issue_data.get('start_date', 'Not set')} |
| **Due date** | {issue_data.get('due_date', 'Not set')} |
| **Estimated hours** | {estimated_text} |
| **Spent hours** | {spent_text} |
| **Actual end date** | {actual_end_date or 'Not set'} |
| **Resolution date** | {resolve_date or 'Not set'} |
| **Updated** | {issue_data.get('updated_on', 'N/A')} |"""

        # Parent issue
        if 'parent' in issue_data and issue_data['parent']:
            parent = issue_data['parent']
            result += f"\n| **Parent Issue** | #{parent['id']} - {parent.get('subject', 'N/A')} |"

        # Children
        if include_details and 'children' in issue_data and issue_data['children']:
            children_text = ", ".join([f"#{c.get('id')}" for c in issue_data['children']])
            result += f"\n| **Sub-topic** | {children_text} |"

        # Custom fields (excluding id=23, 64 that have been displayed independently)
        other_cf = [cf for cf in custom_fields if cf.get('id') not in (23, 64) and cf.get('value', '')]
        if other_cf:
            for cf in other_cf:
                cf_name = cf.get('name', f"ID:{cf.get('id')}")
                result += f"\n| **{cf_name}** (ID:{cf.get('id')}) | {cf.get('value', '')} |"

        # describe
        description = issue_data.get('description', '').strip()
        result += f"\n\n## Description\n\n{description if description else '(no description)'}"

        # appendix
        if include_details and 'attachments' in issue_data and issue_data['attachments']:
            result += f"\n\n## Attachment ({len(issue_data['attachments'])} )\n"
            result += "\n| File name | Size | Type | Uploader | ID |"
            result += "\n|------|------|------|--------|-----|"
            for att in issue_data['attachments']:
                file_size = att.get('filesize', 0)
                size_text = f"{file_size / (1024 * 1024):.2f} MB" if file_size >= 1024 * 1024 else f"{file_size / 1024:.1f} KB"
                result += f"\n| {att.get('filename', 'N/A')} | {size_text} | {att.get('content_type', 'N/A')} | {att.get('author', {}).get('name', 'N/A')} | {att.get('id', '')} |"

        # Notes/History
        if include_details and 'journals' in issue_data and issue_data['journals']:
            notes_journals = [j for j in issue_data['journals'] if j.get('notes', '').strip()]
            if notes_journals:
                result += f"\n\n## Remarks record ({len(notes_journals)} Pen)\n"
                for i, journal in enumerate(notes_journals, 1):
                    author_name = journal.get('user', {}).get('name', 'N/A')
                    created_on = journal.get('created_on', 'N/A')
                    notes = journal.get('notes', '').strip()
                    result += f"\n### #{i} - {author_name}（{created_on}）\n\n{notes}\n"

        # Sub-topic details
        if include_details and 'children' in issue_data and issue_data['children']:
            result += f"\n\n## Sub-topic ({len(issue_data['children'])} )\n"
            result += "\n| ID | Category | Title | Status |"
            result += "\n|----|------|------|------|"
            for child in issue_data['children']:
                child_tracker = child.get('tracker', {}).get('name', '')
                child_status = child.get('status', {}).get('name', 'N/A') if child.get('status') else 'N/A'
                result += f"\n| #{child.get('id')} | {child_tracker} | {child.get('subject', 'N/A')} | {child_status} |"

        # Related issues
        if include_details and 'relations' in issue_data and issue_data['relations']:
            result += f"\n\n## Related issues ({len(issue_data['relations'])} )\n"
            result += "\n| Issue | Association Type |"
            result += "\n|------|----------|"
            for rel in issue_data['relations']:
                rel_type = rel.get('relation_type', 'relates')
                other_id = rel.get('issue_to_id') if rel.get('issue_id') == issue_data['id'] else rel.get('issue_id')
                result += f"\n| #{other_id} | {rel_type} |"

        # observer
        if include_details and 'watchers' in issue_data and issue_data['watchers']:
            watcher_names = [w.get('name', f"ID:{w.get('id')}") for w in issue_data['watchers']]
            result += f"\n\n## Observer ({len(watcher_names)} People)\n\n{', '.join(watcher_names)}"

        return result
        
    except RedmineAPIError as e:
        return f"Failed to get issue: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def update_issue_status(issue_id: int, status_id: int = None, status_name: str = None, notes: str = "") -> str:
    """
    Update issue status
    
    Args:
        issue_id: Issue ID
        status_id: new status ID (optional with status_name)
        status_name: new status name (choose one from status_id)
        notes: update notes (optional)
    
    Returns:
        Update result message
    """
    try:
        client = get_client()
        
        # Processing status parameters
        final_status_id = status_id
        if status_name:
            final_status_id = client.find_status_id_by_name(status_name)
            if not final_status_id:
                return f'Status name not found: "{status_name}"\n\n Available status: \n' + "\n".join([f"- {name}" for name in client.get_available_statuses().keys()])
        
        if not final_status_id:
            return "Error: One of the status_id or status_name parameters must be supplied"
        
        # Prepare to update information
        update_data = {'status_id': final_status_id}
        if notes.strip():
            update_data['notes'] = notes.strip()
        
        # perform update
        client.update_issue(issue_id, **update_data)
        
        # Get confirmation of updated issue information
        updated_issue = client.get_issue(issue_id)
        
        result = f"""Issue status updated successfully!

Issue: #{issue_id} - {updated_issue.subject}
New status: {updated_issue.status.get('name', 'N/A')}"""

        if notes.strip():
            result += f"\nRemarks: {notes}"
            
        return result
        
    except RedmineAPIError as e:
        return f"Failed to update issue status: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def list_project_issues(project_id: int, status_filter: str = "open", limit: int = 20) -> str:
    """
    List project topics
    
    Args:
        project_id: Project ID
        status_filter: status filter ("open", "closed", "all")
        limit: maximum number of return requests (default 20, maximum 100)
    
    Returns:
        List of project issues, presented in table format
    """
    try:
        client = get_client()
        
        # limit limit range
        limit = min(max(limit, 1), 100)
        
        # Filter setting parameters based on status
        params = {
            'project_id': project_id,
            'limit': limit,
            'sort': 'updated_on:desc'
        }
        
        # Processing status filter
        if status_filter == "open":
            params['status_id'] = 'o'  # Redmine API uses 'o' to indicate open status
        elif status_filter == "closed":
            params['status_id'] = 'c'  # Redmine API uses 'c' to indicate closed status
        # "all" does not set status_id
        
        # Get a list of issues
        issues = client.list_issues(**params)
        
        if not issues:
            return f"project {project_id} No matching issues were found in"
        
        # Get project information
        try:
            project = client.get_project(project_id)
            project_name = project.name
        except:
            project_name = f"project {project_id}"
        
        # Format issue list
        result = f"""Project: {project_name}
Status filter: {status_filter}
Found {len(issues)} issues:

{"ID":<8} {"Subject":<40} {"Status":<12} {"Assigned to":<15} {"Updated":<10}
{"-"*8} {"-"*40} {"-"*12} {"-"*15} {"-"*10}"""

        for issue in issues:
            title = issue.subject[:37] + "..." if len(issue.subject) > 40 else issue.subject
            status = issue.status.get('name', 'N/A')[:10]
            assignee = issue.assigned_to.get('name', 'Not assigned')[:13] if issue.assigned_to else 'Not assigned'
            updated = issue.updated_on[:10] if issue.updated_on else 'N/A'
            
            result += f"\n{issue.id:<8} {title:<40} {status:<12} {assignee:<15} {updated:<10}"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to list project issues: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_issue_statuses() -> str:
    """
    Get a list of all available issue statuses
    
    Returns:
        Formatted status list
    """
    try:
        client = get_client()
        statuses = client.get_issue_statuses()
        
        if not statuses:
            return "Issue status not found"
        
        result = "Available issue statuses:\n\n"
        result += f"{'ID':<5} {'name':<15} {'Closed':<8}\n"
        result += f"{'-'*5} {'-'*15} {'-'*8}\n"
        
        for status in statuses:
            is_closed = "yes" if status.get('is_closed', False) else "no"
            result += f"{status['id']:<5} {status['name']:<15} {is_closed:<8}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to get issue status: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_trackers() -> str:
    """
    Get a list of all available trackers
    
    Returns:
        Formatted tracker list
    """
    try:
        client = get_client()
        trackers = client.get_trackers()
        
        if not trackers:
            return "No tracker found"
        
        result = "Available trackers:\n\n"
        result += f"{'ID':<5} {'name':<20} {'Default state':<12}\n"
        result += f"{'-'*5} {'-'*20} {'-'*12}\n"
        
        for tracker in trackers:
            default_status = tracker.get('default_status', {}).get('name', 'N/A')
            result += f"{tracker['id']:<5} {tracker['name']:<20} {default_status:<12}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to get tracker list: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_priorities() -> str:
    """
    Get a priority list of all available issues
    
    Returns:
        Formatted priority list
    """
    try:
        client = get_client()
        priorities = client.get_priorities()
        
        if not priorities:
            return "Issue priority not found"
        
        result = "Available issue priorities:\n\n"
        result += f"{'ID':<5} {'name':<15} {'default':<8}\n"
        result += f"{'-'*5} {'-'*15} {'-'*8}\n"
        
        for priority in priorities:
            is_default = "yes" if priority.get('is_default', False) else "no"
            result += f"{priority['id']:<5} {priority['name']:<15} {is_default:<8}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to obtain issue priority: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_time_entry_activities() -> str:
    """
    Get a list of all available time tracking activities
    
    Returns:
        Formatted time tracking activity list
    """
    try:
        client = get_client()
        activities = client.get_time_entry_activities()
        
        if not activities:
            return "No time tracking activity found"
        
        result = "Available time tracking activities:\n\n"
        result += f"{'ID':<5} {'name':<20} {'default':<8}\n"
        result += f"{'-'*5} {'-'*20} {'-'*8}\n"
        
        for activity in activities:
            is_default = "yes" if activity.get('is_default', False) else "no"
            result += f"{activity['id']:<5} {activity['name']:<20} {is_default:<8}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to get time tracking activity: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_document_categories() -> str:
    """
    Get a list of all available file categories
    
    Returns:
        Formatted file classification list
    """
    try:
        client = get_client()
        categories = client.get_document_categories()
        
        if not categories:
            return "No file category found"
        
        result = "Available file categories:\n\n"
        result += f"{'ID':<5} {'name':<25} {'default':<8}\n"
        result += f"{'-'*5} {'-'*25} {'-'*8}\n"
        
        for category in categories:
            is_default = "yes" if category.get('is_default', False) else "no"
            result += f"{category['id']:<5} {category['name']:<25} {is_default:<8}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to get file classification: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_issue_categories(project_id: int) -> str:
    """
    Obtain the issue classification list of the specified project

    Topic classification is a project-level setting, and each project can have a different classification.

    Args:
        project_id: Project ID

    Returns:
        Formatted issue category list
    """
    try:
        client = get_client()
        categories = client.get_issue_categories(project_id)

        if not categories:
            return f"project #{project_id} No topic classification is set"

        result = f"project #{project_id} Topic classification:\n\n"
        result += f"{'ID':<5} {'name':<25}\n"
        result += f"{'-'*5} {'-'*25}\n"

        for category in categories:
            result += f"{category['id']:<5} {category['name']:<25}\n"

        return result

    except RedmineAPIError as e:
        return f"Failed to obtain issue classification: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_projects() -> str:
    """
    Get a list of accessible projects
    
    Returns:
        Formatted project list
    """
    try:
        client = get_client()
        projects = client.list_projects(limit=50)
        
        if not projects:
            return "No accessible project found"
        
        result = f"Found {len(projects)} projects:\n\n"
        result += f"{'ID':<5} {'Identification code':<20} {'name':<30} {'state':<8}\n"
        result += f"{'-'*5} {'-'*20} {'-'*30} {'-'*8}\n"
        
        for project in projects:
            status_text = "normal" if project.status == 1 else "Seal"
            name = project.name[:27] + "..." if len(project.name) > 30 else project.name
            result += f"{project.id:<5} {project.identifier:<20} {name:<30} {status_text:<8}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to get project list: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def search_issues(query: str, project_id: int = None, limit: int = 10) -> str:
    """
    Search for issues (search keywords in title or description)
    
    Args:
        query: search keyword
        project_id: restrict search to specific projects (optional)
        limit: maximum number of return requests (default 10, maximum 50)
    
    Returns:
        List of issues matching search criteria
    """
    try:
        if not query.strip():
            return "Please provide search keywords"
        
        client = get_client()
        limit = min(max(limit, 1), 50)
        
        # Set search parameters
        params = {
            'limit': limit * 3,  # Get more results to filter
            'sort': 'updated_on:desc'
        }
        
        if project_id:
            params['project_id'] = project_id
        
        # Get a list of issues
        all_issues = client.list_issues(**params)
        
        # Perform keyword filtering locally (because the Redmine API does not have built-in search)
        query_lower = query.lower()
        matching_issues = []
        
        for issue in all_issues:
            if (query_lower in issue.subject.lower() or 
                (issue.description and query_lower in issue.description.lower())):
                matching_issues.append(issue)
                if len(matching_issues) >= limit:
                    break
        
        if not matching_issues:
            search_scope = f"project {project_id}" if project_id else "All accessible projects"
            return f"exist {search_scope} Not found containing '{query}' issues"
        
        # Format results
        result = f"Search keyword: '{query}'\n"
        if project_id:
            result += f"Search scope: Project {project_id}\n"
        result += f"Found {len(matching_issues)} related issues:\n\n"
        
        result += f"{'ID':<8} {'title':<35} {'state':<12} {'project':<15}\n"
        result += f"{'-'*8} {'-'*35} {'-'*12} {'-'*15}\n"
        
        for issue in matching_issues:
            title = issue.subject[:32] + "..." if len(issue.subject) > 35 else issue.subject
            status = issue.status.get('name', 'N/A')[:10]
            project_name = issue.project.get('name', 'N/A')[:13]
            
            result += f"{issue.id:<8} {title:<35} {status:<12} {project_name:<15}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Search for issue failed: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def update_issue_content(issue_id: int, subject: str = None, description: str = None,
                        priority_id: int = None, priority_name: str = None,
                        done_ratio: int = None, tracker_id: int = None, tracker_name: str = None,
                        parent_issue_id: int = None, remove_parent: bool = False, start_date: str = None, due_date: str = None,
                        estimated_hours: float = None,
                        category_id: int = None, category_name: str = None,
                        custom_fields: list[dict] = None) -> str:
    """
    Update issue content (title, description, priority, completion, tracker, category, date, hours, custom fields, etc.)

    Args:
        issue_id: Issue ID
        subject: new topic title (optional)
        description: new issue description (optional)
        priority_id: new priority ID (choose one from priority_name)
        priority_name: new priority name (choose one from priority_id)
        done_ratio: new completion percentage 0-100 (optional)
        tracker_id: new tracker ID (alternative to tracker_name)
        tracker_name: new tracker name (choose one from tracker_id)
        parent_issue_id: new parent issue ID (optional)
        remove_parent: whether to remove the parent issue relationship (optional)
        start_date: new start date in YYYY-MM-DD format (optional)
        due_date: new completion date in YYYY-MM-DD format (optional)
        estimated_hours: new estimated hours (optional)
        category_id: Issue category ID (choose one from category_name, optional, project level)
        category_name: topic category name (choose one from category_id, optional, you need to obtain the topic first to confirm the project it belongs to)
        custom_fields: Custom fields list (optional), format: [{"id": 23, "value": "2026-03-01"}, ...]

    Returns:
        Update result message
    """
    try:
        client = get_client()
        
        # Prepare to update information
        update_data = {}
        changes = []
        
        if subject is not None:
            update_data['subject'] = subject.strip()
            changes.append(f"Subject: {subject}")
        
        if description is not None:
            update_data['description'] = description
            changes.append("Description updated")
        
        # Handle priority parameters
        if priority_name:
            priority_id = client.find_priority_id_by_name(priority_name)
            if not priority_id:
                return f'Priority name not found: "{priority_name}"\n\nAvailable priorities: \n' + "\n".join([f"- {name}" for name in client.get_available_priorities().keys()])
        
        if priority_id is not None:
            update_data['priority_id'] = priority_id
            changes.append(f"Priority ID: {priority_id}")
        
        if done_ratio is not None:
            if not (0 <= done_ratio <= 100):
                return "Error: Complete percentage must be between 0-100"
            update_data['done_ratio'] = done_ratio
            changes.append(f"Done ratio: {done_ratio}%")
        
        # Handling tracker parameters
        if tracker_name:
            tracker_id = client.find_tracker_id_by_name(tracker_name)
            if not tracker_id:
                return f'Tracker name not found: "{tracker_name}"\n\nAvailable trackers: \n' + "\n".join([f"- {name}" for name in client.get_available_trackers().keys()])
        
        if tracker_id is not None:
            update_data['tracker_id'] = tracker_id
            changes.append(f"Tracker ID: {tracker_id}")
        
        if remove_parent:
            update_data['parent_issue_id'] = None
            changes.append("Remove parent issue relationship")
        elif parent_issue_id is not None:
            update_data['parent_issue_id'] = parent_issue_id
            changes.append(f"Parent issue ID: {parent_issue_id}")
        
        if start_date is not None:
            # Validate date format
            try:
                from datetime import datetime
                datetime.strptime(start_date, '%Y-%m-%d')
                update_data['start_date'] = start_date
                changes.append(f"start date: {start_date}")
            except ValueError:
                return "Error: Start date format must be YYYY-MM-DD"
        
        if due_date is not None:
            # Validate date format
            try:
                from datetime import datetime
                datetime.strptime(due_date, '%Y-%m-%d')
                update_data['due_date'] = due_date
                changes.append(f"Completion date: {due_date}")
            except ValueError:
                return "Error: Completion date format must be YYYY-MM-DD"
        
        if estimated_hours is not None:
            if estimated_hours < 0:
                return "Error: Estimated hours cannot be negative"
            update_data['estimated_hours'] = estimated_hours
            changes.append(f"Estimated hours: {estimated_hours}")

        if category_name:
            # You need to get the project_id of the issue first to query the classification
            issue_info = client.get_issue(issue_id)
            project_id = issue_info.project.get('id')
            category_id = client.find_category_id_by_name(project_id, category_name)
            if not category_id:
                available = client.get_available_categories(project_id)
                if available:
                    return f'Category name not found: "{category_name}"\n\nAvailable categories for this project: \n' + "\n".join([f"- {name}" for name in available.keys()])
                else:
                    return f'Category name not found: "{category_name}"\n\nThis project has not yet set an issue classification'

        if category_id is not None:
            update_data['category_id'] = category_id
            changes.append(f"Category ID: {category_id}")

        if custom_fields is not None:
            if not isinstance(custom_fields, list):
                return "Error: custom_fields must be in list format, e.g. [{\"id\": 23, \"value\": \"2026-03-01\"}]"
            for cf in custom_fields:
                if 'id' not in cf or 'value' not in cf:
                    return "Error: Each custom field must contain 'id' and 'value', for example {\"id\": 23, \"value\": \"2026-03-01\"}"
            update_data['custom_fields'] = custom_fields
            changes.append(f"Custom fields: {len(custom_fields)} fields updated")

        if not update_data and not changes:
            return "Error: Please provide at least one field to update"
        
        # perform update
        client.update_issue(issue_id, **update_data)
        
        # Get updated issue information
        updated_issue = client.get_issue(issue_id)
        
        result = f"""Issue content updated successfully!

Issue: #{issue_id} - {updated_issue.subject}
Updated fields:
{chr(10).join(f"- {change}" for change in changes)}

Current status:
- Tracker: {updated_issue.tracker.get('name', 'N/A')}
- Status: {updated_issue.status.get('name', 'N/A')}
- Priority: {updated_issue.priority.get('name', 'N/A')}
- Done ratio: {updated_issue.done_ratio}%"""

        return result
        
    except RedmineAPIError as e:
        return f"Failed to update issue content: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def add_issue_note(issue_id: int, notes: str, private: bool = False, 
                   spent_hours: float = None, activity_name: str = None, 
                   activity_id: int = None, spent_on: str = None) -> str:
    """
    Add notes to the topic and record the time at the same time
    
    Args:
        issue_id: Issue ID
        notes: Note content
        private: whether the comment is private (default is no)
        spent_hours: Spent working hours (hours)
        activity_name: activity name (choose one from activity_id)
        activity_id: Activity ID (choose one from activity_name)
        spent_on: Record date in YYYY-MM-DD format (optional, default is today)
    
    Returns:
        Add result message
    """
    try:
        if not notes.strip():
            return "Error: Note content cannot be empty"
        
        client = get_client()
        time_entry_id = None
        
        # Processing time record
        if spent_hours is not None:
            if spent_hours <= 0:
                return "Error: Spent hours must be greater than 0"
            
            # Handle activity parameters
            final_activity_id = activity_id
            if activity_name:
                final_activity_id = client.find_time_entry_activity_id_by_name(activity_name)
                if not final_activity_id:
                    available_activities = client.get_available_time_entry_activities()
                    return f'Time tracking activity name not found: "{activity_name}"\n\nAvailable activities: \n' + "\n".join([f"- {name}" for name in available_activities.keys()])
            
            if not final_activity_id:
                return "Error: activity_id or activity_name parameter must be provided"
            
            # Create time record
            try:
                time_entry_id = client.create_time_entry(
                    issue_id=issue_id,
                    hours=spent_hours,
                    activity_id=final_activity_id,
                    comments=notes.strip(),
                    spent_on=spent_on
                )
            except Exception as e:
                return f"Failed to create time record: {str(e)}"
        
        # Prepare updated information (added remarks)
        update_data = {'notes': notes.strip()}
        if private:
            update_data['private_notes'] = True
        
        # perform update
        client.update_issue(issue_id, **update_data)
        
        # Get issue information
        issue = client.get_issue(issue_id)
        
        privacy_text = "private" if private else "public"
        result = f"""Note added successfully!

Issue: #{issue_id} - {issue.subject}
Note type: {privacy_text}
Remarks:
{notes.strip()}"""

        # If there is a creation time record, add relevant information
        if time_entry_id:
            from datetime import date
            actual_date = spent_on if spent_on else date.today().strftime('%Y-%m-%d')
            activity_name_display = activity_name if activity_name else f"ID {final_activity_id}"
            result += f"""

Time record added successfully!
- Time record ID: {time_entry_id}
- Man hours consumed: {spent_hours} hours
- Activities: {activity_name_display}
- Date recorded: {actual_date}"""

        return result
        
    except RedmineAPIError as e:
        return f"Failed to add topic notes: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def assign_issue(issue_id: int, user_id: int = None, user_name: str = None, user_login: str = None, notes: str = "") -> str:
    """
    Assign issues to users
    
    Args:
        issue_id: Issue ID
        user_id: assigned user ID (choose one from user_name/user_login)
        user_name: assigned user name (choose one from user_id/user_login)
        user_login: assigned user login name (choose one from user_id/user_name)
        notes: Assign notes (optional)
    
    Returns:
        Assignment result message
    """
    try:
        client = get_client()
        
        # Handle user parameters
        final_user_id = user_id
        if user_name:
            final_user_id = client.find_user_id_by_name(user_name)
            if not final_user_id:
                users = client.get_available_users()
                return f'Username not found: "{user_name}"\n\n Available users (name): \n' + "\n".join([f"- {name}" for name in users['by_name'].keys()])
        elif user_login:
            final_user_id = client.find_user_id_by_login(user_login)
            if not final_user_id:
                users = client.get_available_users()
                return f'User login name not found: "{user_login}"\n\n Available users (login name): \n' + "\n".join([f"- {login}" for login in users['by_login'].keys()])
        
        # Prepare to update information
        update_data = {}
        
        if final_user_id is not None:
            update_data['assigned_to_id'] = final_user_id
            action_text = f"Assigned to user ID {final_user_id}"
        else:
            update_data['assigned_to_id'] = None
            action_text = "Unassign"
        
        if notes.strip():
            update_data['notes'] = notes.strip()
        
        # perform update
        client.update_issue(issue_id, **update_data)
        
        # Get updated issue information
        updated_issue = client.get_issue(issue_id)
        
        assignee_name = "Not assigned"
        if updated_issue.assigned_to:
            assignee_name = updated_issue.assigned_to.get('name', f"User ID {user_id}")
        
        result = f"""Topic assignment updated successfully!

Issue: #{issue_id} - {updated_issue.subject}
Action: {action_text}
Currently assigned to: {assignee_name}"""

        if notes.strip():
            result += f"\nRemarks: {notes}"

        return result
        
    except RedmineAPIError as e:
        return f"Failed to assign issue: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def create_new_issue(project_id: int, subject: str, description: str = "",
                    tracker_id: int = None, tracker_name: str = None,
                    priority_id: int = None, priority_name: str = None,
                    assigned_to_id: int = None, assigned_to_name: str = None, assigned_to_login: str = None,
                    parent_issue_id: int = None, start_date: str = None, due_date: str = None,
                    estimated_hours: float = None, status_id: int = None, status_name: str = None,
                    category_id: int = None, category_name: str = None) -> str:
    """
    Create a new Redmine issue

    Recommended title format: [Type] Problem description, for example:
    - [Function Development] Evaluation submission completion page logic
    - [Interface Correction] The question list switches the language and the Dialog text is unified
    - [Backend Bug] The default expansion and delete buttons for new questions are abnormal

    Commonly used types: [Interface modification] [Function development] [Front-end optimization] [Backend Bug] [E2E testing] [Database] [Devops] [Refactoring] [Documents]

    Args:
        project_id: Project ID
        subject: topic title (it is recommended to add the [type] prefix)
        description: topic description (optional)
        tracker_id: Tracker ID (optional with tracker_name)
        tracker_name: tracker name (choose one from tracker_id)
        priority_id: priority ID (choose one from priority_name)
        priority_name: priority name (choose one from priority_id)
        assigned_to_id: assigned user ID (choose one from assigned_to_name/assigned_to_login)
        assigned_to_name: assigned user name (choose one of three assigned_to_id/assigned_to_login)
        assigned_to_login: assigned user login name (choose one of three assigned_to_id/assigned_to_name)
        parent_issue_id: parent issue ID (optional, used to create child issues)
        start_date: start date in YYYY-MM-DD format (optional)
        due_date: completion date in YYYY-MM-DD format (optional)
        estimated_hours: estimated hours (optional)
        status_id: initial status ID (choose one from status_name, optional)
        status_name: initial status name (choose one from status_id, optional)
        category_id: Issue category ID (choose one from category_name, optional, project level)
        category_name: Issue category name (choose one from category_id, optional, project level)

    Returns:
        Create result message
    """
    try:
        if not subject.strip():
            return "Error: Issue title cannot be empty"
        
        client = get_client()
        
        # Handling tracker parameters
        final_tracker_id = tracker_id
        if tracker_name:
            final_tracker_id = client.find_tracker_id_by_name(tracker_name)
            if not final_tracker_id:
                return f'Tracker name not found: "{tracker_name}"\n\nAvailable trackers: \n' + "\n".join([f"- {name}" for name in client.get_available_trackers().keys()])
        
        # Handle priority parameters
        final_priority_id = priority_id
        if priority_name:
            final_priority_id = client.find_priority_id_by_name(priority_name)
            if not final_priority_id:
                return f'Priority name not found: "{priority_name}"\n\nAvailable priorities: \n' + "\n".join([f"- {name}" for name in client.get_available_priorities().keys()])
        
        # Handling assigned user parameters
        final_assigned_to_id = assigned_to_id
        if assigned_to_name:
            final_assigned_to_id = client.find_user_id_by_name(assigned_to_name)
            if not final_assigned_to_id:
                users = client.get_available_users()
                return f'Username not found: "{assigned_to_name}"\n\n Available users (name): \n' + "\n".join([f"- {name}" for name in users['by_name'].keys()])
        elif assigned_to_login:
            final_assigned_to_id = client.find_user_id_by_login(assigned_to_login)
            if not final_assigned_to_id:
                users = client.get_available_users()
                return f'User login name not found: "{assigned_to_login}"\n\n Available users (login name): \n' + "\n".join([f"- {login}" for login in users['by_login'].keys()])
        
        # Processing status parameters
        final_status_id = status_id
        if status_name:
            final_status_id = client.find_status_id_by_name(status_name)
            if not final_status_id:
                return f'Status name not found: "{status_name}"\n\n Available status: \n' + "\n".join([f"- {name}" for name in client.get_available_statuses().keys()])

        # Handle classification parameters
        final_category_id = category_id
        if category_name:
            final_category_id = client.find_category_id_by_name(project_id, category_name)
            if not final_category_id:
                available = client.get_available_categories(project_id)
                if available:
                    return f'Category name not found: "{category_name}"\n\nAvailable categories for this project: \n' + "\n".join([f"- {name}" for name in available.keys()])
                else:
                    return f'Category name not found: "{category_name}"\n\nThis project has not yet set an issue classification'

        # Validate date format
        for date_label, date_val in [("start date", start_date), ("completion date", due_date)]:
            if date_val is not None:
                try:
                    datetime.strptime(date_val, '%Y-%m-%d')
                except ValueError:
                    return f"mistake: {date_label}The format must be YYYY-MM-DD"

        if estimated_hours is not None and estimated_hours < 0:
            return "Error: Estimated hours cannot be negative"

        # Create an issue
        new_issue_id = client.create_issue(
            project_id=project_id,
            subject=subject.strip(),
            description=description,
            tracker_id=final_tracker_id,
            priority_id=final_priority_id,
            assigned_to_id=final_assigned_to_id,
            parent_issue_id=parent_issue_id,
            status_id=final_status_id,
            start_date=start_date,
            due_date=due_date,
            estimated_hours=estimated_hours,
            category_id=final_category_id,
        )
        
        # Get created issue information
        new_issue = client.get_issue(new_issue_id)
        
        result = f"""New issue created successfully!

Issue ID: #{new_issue_id}
Subject: {new_issue.subject}
Project: {new_issue.project.get('name', 'N/A')}
Tracker: {new_issue.tracker.get('name', 'N/A')}
Status: {new_issue.status.get('name', 'N/A')}
Priority: {new_issue.priority.get('name', 'N/A')}
Assigned to: {new_issue.assigned_to.get('name', 'Not assigned') if new_issue.assigned_to else 'Not assigned'}"""

        if parent_issue_id:
            result += f"\nParent issue: #{parent_issue_id}"
        if start_date:
            result += f"\nStart date: {start_date}"
        if due_date:
            result += f"\nCompletion date: {due_date}"
        if estimated_hours is not None:
            result += f"\nEstimated hours: {estimated_hours}"

        if description:
            result += f"\n\nDescription:\n{description}"

        return result
        
    except RedmineAPIError as e:
        return f"Failed to create issue: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_my_issues(status_filter: str = "open", limit: int = 20) -> str:
    """
    Get a list of issues assigned to me
    
    Args:
        status_filter: status filter ("open", "closed", "all")
        limit: maximum number of return requests (default 20, maximum 100)
    
    Returns:
        My issue list
    """
    try:
        client = get_client()
        
        # Get current user information first
        current_user = client.get_current_user()
        user_id = current_user['id']
        user_name = current_user.get('firstname', '') + ' ' + current_user.get('lastname', '')
        
        # limit limit range
        limit = min(max(limit, 1), 100)
        
        # Set query parameters
        params = {
            'assigned_to_id': user_id,
            'limit': limit,
            'sort': 'updated_on:desc'
        }
        
        # Processing status filter
        if status_filter == "open":
            params['status_id'] = 'o'  # Redmine API uses 'o' to indicate open status
        elif status_filter == "closed":
            params['status_id'] = 'c'  # Redmine API uses 'c' to indicate closed status
        
        # Get a list of issues
        issues = client.list_issues(**params)
        
        if not issues:
            return f"Assigned to {user_name.strip()} of{status_filter}issue"
        
        # Format results
        result = f"""Issues assigned to {user_name.strip()}:
Status filter: {status_filter}
Found {len(issues)} issues:

{"ID":<8} {"title":<35} {"project":<15} {"state":<12} {"Update time":<10}
{"-"*8} {"-"*35} {"-"*15} {"-"*12} {"-"*10}"""

        for issue in issues:
            title = issue.subject[:32] + "..." if len(issue.subject) > 35 else issue.subject
            project_name = issue.project.get('name', 'N/A')[:13]
            status = issue.status.get('name', 'N/A')[:10]
            updated = issue.updated_on[:10] if issue.updated_on else 'N/A'
            
            result += f"\n{issue.id:<8} {title:<35} {project_name:<15} {status:<12} {updated:<10}"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to get my issue: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def close_issue(issue_id: int, notes: str = "", done_ratio: int = 100) -> str:
    """
    Close the issue (set to completed status)
    
    Args:
        issue_id: Issue ID
        notes: Close notes (optional)
        done_ratio: completion percentage (default 100%)
    
    Returns:
        Close results message
    """
    try:
        client = get_client()
        
        # Get a list of available states and look for closed states
        statuses = client.get_issue_statuses()
        closed_status_id = None
        
        for status in statuses:
            if status.get('is_closed', False):
                closed_status_id = status['id']
                break
        
        if closed_status_id is None:
            return "Error: No available shutdown status found"
        
        # Prepare to update information
        update_data = {
            'status_id': closed_status_id,
            'done_ratio': min(max(done_ratio, 0), 100)
        }
        
        if notes.strip():
            update_data['notes'] = notes.strip()
        
        # perform update
        client.update_issue(issue_id, **update_data)
        
        # Get updated issue information
        updated_issue = client.get_issue(issue_id)
        
        result = f"""Issue closed successfully!

Issue: #{issue_id} - {updated_issue.subject}
Status: {updated_issue.status.get('name', 'N/A')}
Done ratio: {updated_issue.done_ratio}%"""

        if notes.strip():
            result += f"\nClose Notes: {notes}"

        return result
        
    except RedmineAPIError as e:
        return f"Failed to close issue: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def resolve_issue(
    issue_id: int,
    notes: str = "",
    custom_fields: list[dict] = None,
    status_name: str = "Resolved",
    status_id: int = None,
    done_ratio: int = 100,
    spent_hours: float = None,
    activity_name: str = None,
    activity_id: int = None,
) -> str:
    """
    Mark issue as resolved (composite operation)

    Complete it all at once: update status, set completion level, fill in custom fields, add notes, and optional time records.
    It applies to the standard closing process after development is completed.

    Typical usage:
        resolve_issue(
            issue_id=123,
            notes="commit abc1234\\nModification summary: \\n- Complete functional development\\n- Pass unit testing",
            custom_fields=[
                {"id": 23, "value": "2026-03-07"},
                {"id": 64, "value": "2026-03-07"}
            ]
        )

    Args:
        issue_id: Issue ID
        notes: Solution notes (it is recommended to include commit hash and modification summary)
        custom_fields: custom fields list (optional), for example [{"id": 23, "value": "2026-03-07"}]
        status_name: target status name (default "resolved", can be changed to other status)
        status_id: target status ID (takes precedence over status_name)
        done_ratio: completion percentage (default 100)
        spent_hours: hours spent (optional, record time at the same time)
        activity_name: Time recording activity name (choose one from activity_id, it will take effect only when spent_hours has a value)
        activity_id: Time recording activity ID (choose one from activity_name)

    Returns:
        Resolution result message
    """
    try:
        client = get_client()

        # parsing status
        final_status_id = status_id
        if not final_status_id:
            final_status_id = client.find_status_id_by_name(status_name)
            if not final_status_id:
                return f'Status name not found: "{status_name}"\n\n Available status: \n' + "\n".join([f"- {name}" for name in client.get_available_statuses().keys()])

        # Prepare to update information
        update_data = {
            'status_id': final_status_id,
            'done_ratio': min(max(done_ratio, 0), 100),
        }

        if notes.strip():
            update_data['notes'] = notes.strip()

        if custom_fields:
            for cf in custom_fields:
                if 'id' not in cf or 'value' not in cf:
                    return "Error: Each custom field must contain 'id' and 'value'"
            update_data['custom_fields'] = custom_fields

        # One API call for all updates
        client.update_issue(issue_id, **update_data)

        # Record hours worked (if specified)
        if spent_hours is not None and spent_hours > 0:
            final_activity_id = activity_id
            if activity_name:
                final_activity_id = client.find_time_entry_activity_id_by_name(activity_name)
                if not final_activity_id:
                    return f'Issue resolved, but time recording failed: Activity name not found"{activity_name}"\n\nAvailable activities: \n' + "\n".join([f"- {name}" for name in client.get_available_time_entry_activities().keys()])

            if final_activity_id:
                today = datetime.now().strftime('%Y-%m-%d')
                client.create_time_entry(
                    issue_id=issue_id,
                    hours=spent_hours,
                    activity_id=final_activity_id,
                    comments=notes.strip()[:255] if notes else "",
                    spent_on=today,
                )

        # Get updated issue information
        updated_issue = client.get_issue(issue_id)

        result = f"""Issue resolved successfully!

Issue: #{issue_id} - {updated_issue.subject}
Status: {updated_issue.status.get('name', 'N/A')}
Done ratio: {updated_issue.done_ratio}%"""

        if custom_fields:
            result += f"\nCustom fields: {len(custom_fields)} updated"
        if notes.strip():
            result += f"\n Note: Added"
        if spent_hours is not None and spent_hours > 0:
            result += f"\nTime entry: {spent_hours} hours"

        return result

    except RedmineAPIError as e:
        return f"Issue resolution failed: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def start_working(
    issue_id: int,
    status_name: str = "In practice",
    status_id: int = None,
    notes: str = "",
    set_start_date: bool = True,
    assign_to_me: bool = True,
) -> str:
    """
    Start processing an issue (composite operation)

    Complete it once: Update the status to Implementing, set the start date, and assign it to yourself.
    Form a complete workflow closed loop with resolve_issue.

    Args:
        issue_id: Issue ID
        status_name: target status name (default is "implementing")
        status_id: target status ID (takes precedence over status_name)
        notes: Notes to start processing (optional)
        set_start_date: Whether to automatically set the start date to today (default True)
        assign_to_me: Whether to assign it to yourself (default True, use the user corresponding to the API Key)

    Returns:
        Operation result message
    """
    try:
        client = get_client()

        # parsing status
        final_status_id = status_id
        if not final_status_id:
            final_status_id = client.find_status_id_by_name(status_name)
            if not final_status_id:
                return f'Status name not found: "{status_name}"\n\n Available status: \n' + "\n".join([f"- {name}" for name in client.get_available_statuses().keys()])

        update_data = {
            'status_id': final_status_id,
        }

        if set_start_date:
            update_data['start_date'] = datetime.now().strftime('%Y-%m-%d')

        if assign_to_me:
            current_user = client.get_current_user()
            update_data['assigned_to_id'] = current_user['id']

        if notes.strip():
            update_data['notes'] = notes.strip()

        client.update_issue(issue_id, **update_data)

        # Read updated issues (trigger snapshot saving at the same time)
        include_params = ['journals', 'attachments', 'children', 'relations', 'watchers']
        issue_data = client.get_issue_raw(issue_id, include=include_params)

        assigned_name = issue_data.get('assigned_to', {}).get('name', 'Not assigned') if issue_data.get('assigned_to') else 'Not assigned'

        result = f"""Start working on the issue!

Issue: #{issue_id} - {issue_data['subject']}
Status: {issue_data['status'].get('name', 'N/A')}
Assigned to: {assigned_name}
start date: {issue_data.get('start_date', 'N/A')}"""

        return result

    except RedmineAPIError as e:
        return f"Failed to start processing issue: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def check_issue_changes(issue_id: int) -> str:
    """
    Detect changes to an issue since the last time it was read

    Compare the current status of the issue with the snapshot when it was last read, and list all changes.
    You need to use get_issue to read the issue before the snapshot is available for comparison.

    Args:
        issue_id: Issue ID

    Returns:
        Summary of changes or no change notification
    """
    try:
        client = get_client()
        snapshot = client.get_issue_snapshot(issue_id)

        if not snapshot:
            return f"Issue #{issue_id} There is no snapshot record, please use get_issue( first{issue_id}) read issue"

        # Retrieve the latest information on the topic
        include_params = ['journals', 'attachments', 'children', 'relations', 'watchers']
        current = client.get_issue_raw(issue_id, include=include_params)

        changes = []

        # Compare basic fields
        current_status = current.get('status', {}).get('name')
        if current_status != snapshot['status']:
            changes.append(f"Status: {snapshot['status']} → {current_status}")

        current_priority = current.get('priority', {}).get('name')
        if current_priority != snapshot['priority']:
            changes.append(f"Priority: {snapshot['priority']} → {current_priority}")

        current_assigned = current.get('assigned_to', {}).get('name') if current.get('assigned_to') else None
        if current_assigned != snapshot['assigned_to']:
            changes.append(f"Assigned to: {snapshot['assigned_to'] or 'Not assigned'} → {current_assigned or 'Not assigned'}")

        current_done = current.get('done_ratio', 0)
        if current_done != snapshot['done_ratio']:
            changes.append(f"Done ratio: {snapshot['done_ratio']}% → {current_done}%")

        current_subject = current.get('subject')
        if current_subject != snapshot['subject']:
            changes.append(f"Subject: {snapshot['subject']} → {current_subject}")

        # Describe the differences
        current_desc = current.get('description', '')
        old_desc = snapshot.get('description', '')
        if current_desc != old_desc:
            # Calculate difference summary
            old_lines = old_desc.splitlines() if old_desc else []
            new_lines = current_desc.splitlines() if current_desc else []
            added = len(new_lines) - len(old_lines)
            if added > 0:
                changes.append(f"Description: New contract {added} OK")
            elif added < 0:
                changes.append(f"Description: Reduce approx. {abs(added)} OK")
            else:
                changes.append("Description: Content has been modified")

            # Show specific diff
            import difflib
            diff = list(difflib.unified_diff(
                old_lines, new_lines,
                fromfile='last read', tofile='at present',
                lineterm='', n=2
            ))
            if diff:
                diff_text = "\n".join(diff[:20])  # Display up to 20 lines
                if len(diff) > 20:
                    diff_text += f"\n... (Total {len(diff)} row differences)"
                changes.append(f"Description difference:\n{diff_text}")

        # Remarks changes
        current_journals = current.get('journals', [])
        if len(current_journals) > snapshot['journals_count']:
            new_count = len(current_journals) - snapshot['journals_count']
            changes.append(f"New {new_count} Note:")
            for j in current_journals[-new_count:]:
                author = j.get('user', {}).get('name', 'N/A')
                notes = j.get('notes', '').strip()
                if notes:
                    preview = notes[:100] + '...' if len(notes) > 100 else notes
                    changes.append(f"  - {author}: {preview}")

        # Attachment changes
        current_attachments = current.get('attachments', [])
        if len(current_attachments) > snapshot['attachments_count']:
            new_count = len(current_attachments) - snapshot['attachments_count']
            changes.append(f"New {new_count} Attachments:")
            for att in current_attachments[-new_count:]:
                changes.append(f"  - {att.get('filename', 'N/A')} ({att.get('author', {}).get('name', 'N/A')})")

        # Custom field changes
        current_cf = {
            cf.get('id'): cf.get('value')
            for cf in current.get('custom_fields', [])
        }
        for cf_id, old_val in snapshot.get('custom_fields', {}).items():
            new_val = current_cf.get(cf_id)
            if new_val != old_val:
                # Find the field name
                cf_name = next(
                    (cf.get('name') for cf in current.get('custom_fields', []) if cf.get('id') == cf_id),
                    f"ID:{cf_id}"
                )
                changes.append(f"Custom fields [{cf_name}]: {old_val or '(null)'} → {new_val or '(null)'}")

        # Output results
        snapshot_time = snapshot.get('snapshot_time', 'N/A')

        if not changes:
            return f"Issue #{issue_id} since last read ({snapshot_time}) No changes since"

        result = f"Issue #{issue_id} changes (from {snapshot_time} since):\n"
        result += "=" * 50 + "\n"
        result += "\n".join(changes)

        return result

    except RedmineAPIError as e:
        return f"Failed to detect changes: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def sync_my_issues(project_id: int = None, status_filter: str = "open", limit: int = 20) -> str:
    """
    Batch detection of changes to topics assigned to me

    Get all issues assigned to me and compare with local snapshot,
    Lists a summary of topics that have changed. Issues without snapshots are marked as "first read".

    Args:
        project_id: Project ID (optional, only the issues of this project will be synchronized after specifying)
        status_filter: status filter ("open", "closed", "all")
        limit: maximum number of postbacks (default 20)

    Returns:
        Change summaries for all issues assigned to me
    """
    try:
        client = get_client()

        # Get current user ID
        current_user = client.get_current_user()
        my_user_id = current_user['id']

        # Get my issue list
        status_map = {'open': 'o', 'closed': 'c', 'all': '*', 'o': 'o', 'c': 'c', '*': '*'}
        status_param = status_map.get(status_filter, 'o')
        issues = client.list_issues(
            project_id=project_id,
            assigned_to_id=my_user_id,
            status_id=status_param,
            limit=min(limit, 100),
            sort='updated_on:desc'
        )

        project_hint = f"(project #{project_id}）" if project_id else ""
        if not issues:
            return f"There are currently no issues assigned to me{project_hint}"

        changed_issues = []
        new_issues = []
        unchanged_count = 0

        for issue in issues:
            issue_id = issue.id
            snapshot = client.get_issue_snapshot(issue_id)

            if not snapshot:
                new_issues.append(issue)
                # Read complete data to create a snapshot
                try:
                    client.get_issue_raw(issue_id, include=['journals', 'attachments'])
                except Exception:
                    pass
                continue

            # Quick comparison: whether updated_on is the same
            current_updated = issue.updated_on
            if current_updated == snapshot.get('updated_on'):
                unchanged_count += 1
                continue

            # There are changes, collect summary
            change_summary = []
            current_status = issue.status.get('name')
            if current_status != snapshot['status']:
                change_summary.append(f"Status: {snapshot['status']} → {current_status}")
            if issue.done_ratio != snapshot['done_ratio']:
                change_summary.append(f"Done ratio: {snapshot['done_ratio']}% → {issue.done_ratio}%")

            current_assigned = issue.assigned_to.get('name') if issue.assigned_to else None
            if current_assigned != snapshot['assigned_to']:
                change_summary.append(f"Assignment: {snapshot['assigned_to'] or 'Not assigned'} → {current_assigned or 'Not assigned'}")

            if not change_summary:
                change_summary.append("Content has been updated (please use check_issue_changes for details)")

            changed_issues.append((issue, change_summary))

            # Update snapshot
            try:
                client.get_issue_raw(issue_id, include=['journals', 'attachments'])
            except Exception:
                pass

        # Assembly output
        result = f"My issue sync summary{project_hint}(common {len(issues)} )\n"
        result += "=" * 50 + "\n"

        if changed_issues:
            result += f"\n Topics with changes ({len(changed_issues)} ):\n"
            for issue, summaries in changed_issues:
                title = issue.subject[:40] + '...' if len(issue.subject) > 40 else issue.subject
                result += f"\n  #{issue.id} {title}\n"
                for s in summaries:
                    result += f"    - {s}\n"

        if new_issues:
            result += f"\nThe issue read for the first time ({len(new_issues)} ):\n"
            for issue in new_issues:
                title = issue.subject[:40] + '...' if len(issue.subject) > 40 else issue.subject
                status = issue.status.get('name', 'N/A')
                result += f"  #{issue.id} [{status}] {title}\n"

        if unchanged_count > 0:
            result += f"\n No changes: {unchanged_count} topics"

        return result

    except RedmineAPIError as e:
        return f"Sync failed: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def sync_project_issues(project_id: int, status_filter: str = "open") -> str:
    """
    Synchronize project issue structure and detect changes

    Obtain all issues of the specified project (automatic paging) and present the parent-child hierarchical relationship in a tree structure.
    And compare it with the local snapshot to indicate the change status. Can be executed repeatedly to continuously track project progress.

    Args:
        project_id: project ID (required)
        status_filter: status filter ("open", "closed", "all"), default "open"

    Returns:
        Project topic tree structure and change summary
    """
    try:
        client = get_client()

        # Get all project issues in pages
        status_map = {'open': 'o', 'closed': 'c', 'all': '*', 'o': 'o', 'c': 'c', '*': '*'}
        status_param = status_map.get(status_filter, 'o')
        issues = []
        offset = 0
        page_size = 100  # Redmine API single limit

        while True:
            batch = client.list_issues(
                project_id=project_id,
                status_id=status_param,
                limit=page_size,
                offset=offset,
                sort='updated_on:desc'
            )
            if not batch:
                break
            issues.extend(batch)
            if len(batch) < page_size:
                break
            offset += page_size

        if not issues:
            return f"project #{project_id} There are no matching issues (filter: {status_filter}）"

        # Classification Topics: Changed / Read for the first time / No changes
        changed_issues = []
        new_issues = []
        unchanged_issues = []

        for issue in issues:
            issue_id = issue.id
            snapshot = client.get_issue_snapshot(issue_id)

            if not snapshot:
                new_issues.append(issue)
                # Create a snapshot
                try:
                    client.get_issue_raw(issue_id, include=['journals', 'attachments'])
                except Exception:
                    pass
                continue

            # Quick comparison updated_on
            if issue.updated_on == snapshot.get('updated_on'):
                unchanged_issues.append(issue)
                continue

            # Detect specific changes
            change_summary = []
            current_status = issue.status.get('name')
            if current_status != snapshot['status']:
                change_summary.append(f"Status: {snapshot['status']} → {current_status}")
            if issue.done_ratio != snapshot['done_ratio']:
                change_summary.append(f"Done ratio: {snapshot['done_ratio']}% → {issue.done_ratio}%")
            current_assigned = issue.assigned_to.get('name') if issue.assigned_to else None
            if current_assigned != snapshot['assigned_to']:
                change_summary.append(f"Assignment: {snapshot['assigned_to'] or 'Not assigned'} → {current_assigned or 'Not assigned'}")
            current_priority = issue.priority.get('name')
            if current_priority != snapshot['priority']:
                change_summary.append(f"Priority: {snapshot['priority']} → {current_priority}")
            if issue.subject != snapshot['subject']:
                change_summary.append(f"Title changed")

            if not change_summary:
                change_summary.append("Content has been updated (please use check_issue_changes for details)")

            changed_issues.append((issue, change_summary))

            # Update snapshot
            try:
                client.get_issue_raw(issue_id, include=['journals', 'attachments'])
            except Exception:
                pass

        # Create a parent-child relationship index
        all_issues = {issue.id: issue for issue in issues}
        children_map = {}  # parent_id -> [child_issues]
        root_issues = []

        for issue in issues:
            parent_id = issue.parent.get('id') if hasattr(issue, 'parent') and issue.parent else None
            if parent_id and parent_id in all_issues:
                children_map.setdefault(parent_id, []).append(issue)
            else:
                root_issues.append(issue)

        # Change status query form
        changed_ids = {issue.id for issue, _ in changed_issues}
        new_ids = {issue.id for issue in new_issues}
        change_details = {issue.id: summaries for issue, summaries in changed_issues}

        def get_marker(issue_id):
            if issue_id in changed_ids:
                return "🔄"
            elif issue_id in new_ids:
                return "🆕"
            return "  "

        def format_issue_line(issue, prefix=""):
            marker = get_marker(issue.id)
            tracker = issue.tracker.get('name', '') if hasattr(issue, 'tracker') and issue.tracker else ''
            status = issue.status.get('name', '') if issue.status else ''
            assigned = issue.assigned_to.get('name', 'Not assigned') if issue.assigned_to else 'Not assigned'
            done = issue.done_ratio if hasattr(issue, 'done_ratio') else 0
            title = issue.subject[:50] + '...' if len(issue.subject) > 50 else issue.subject
            return f"{prefix}{marker} #{issue.id} [{tracker}][{status}] {title} ({assigned}, {done}%)"

        def build_tree(issue, prefix="", is_last=True):
            lines = []
            connector = "└── " if is_last else "├── "
            lines.append(format_issue_line(issue, prefix + connector if prefix else ""))

            # Change details
            if issue.id in change_details:
                detail_prefix = prefix + ("    " if is_last else "│   ") if prefix else "    "
                for s in change_details[issue.id]:
                    lines.append(f"{detail_prefix}↳ {s}")

            children = children_map.get(issue.id, [])
            child_prefix = prefix + ("    " if is_last else "│   ") if prefix else ""
            for i, child in enumerate(children):
                is_child_last = (i == len(children) - 1)
                lines.extend(build_tree(child, child_prefix, is_child_last))
            return lines

        # Get project name
        project_name = ""
        if issues:
            project_name = issues[0].project.get('name', '') if hasattr(issues[0], 'project') and issues[0].project else ''

        # Assembly output
        result = f"# Project topic synchronization: {project_name} (#{project_id})\n"
        result += f"filter: {status_filter} | Total {len(issues)} Topics\n"
        result += f"🔄 There are changes: {len(changed_issues)} | 🆕 First read: {len(new_issues)} | ✅ No changes: {len(unchanged_issues)}\n"
        result += "=" * 60 + "\n\n"

        # tree structure
        result += "## Issue structure\n\n"
        for i, issue in enumerate(root_issues):
            is_last = (i == len(root_issues) - 1)
            tree_lines = build_tree(issue, "", is_last)
            result += "\n".join(tree_lines) + "\n"

        # Summary of changes
        if changed_issues:
            result += f"\n## Summary of changes ({len(changed_issues)} )\n\n"
            result += "| Topic | Category | Change content |\n"
            result += "|------|------|----------|\n"
            for issue, summaries in changed_issues:
                tracker = issue.tracker.get('name', '') if hasattr(issue, 'tracker') and issue.tracker else ''
                changes_text = "; ".join(summaries)
                title = issue.subject[:30] + '...' if len(issue.subject) > 30 else issue.subject
                result += f"| #{issue.id} {title} | {tracker} | {changes_text} |\n"

        return result

    except RedmineAPIError as e:
        return f"Failed to synchronize project issues: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def search_users(query: str, limit: int = 10) -> str:
    """
    Search for users (by name or login name)
    
    Args:
        query: search keyword (name or login name)
        limit: maximum number of return requests (default 10, maximum 50)
    
    Returns:
        List of users matching search criteria
    """
    try:
        if not query.strip():
            return "Please provide search keywords"
        
        client = get_client()
        limit = min(max(limit, 1), 50)
        
        users = client.search_users(query, limit)
        
        if not users:
            return f"No users found matching '{query}'"
        
        result = f"Search keyword: '{query}'\n found {len(users)} related users:\n\n"
        result += f"{'ID':<5} {'Login name':<15} {'Name':<20} {'state':<8}\n"
        result += f"{'-'*5} {'-'*15} {'-'*20} {'-'*8}\n"
        
        for user in users:
            full_name = f"{user.firstname} {user.lastname}".strip()
            if not full_name:
                full_name = user.login
            status_text = "enable" if user.status == 1 else "deactivate"
            result += f"{user.id:<5} {user.login:<15} {full_name:<20} {status_text:<8}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to search for user: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def list_users(limit: int = 20, status_filter: str = "active") -> str:
    """
    List all users
    
    Args:
        limit: maximum number of return requests (default 20, maximum 100)
        status_filter: status filter ("active", "locked", "all")
    
    Returns:
        List of users, presented in table format
    """
    try:
        client = get_client()
        limit = min(max(limit, 1), 100)
        
        # Transition status filter
        status = None
        if status_filter == "active":
            status = 1
        elif status_filter == "locked":
            status = 3
        
        users = client.list_users(limit=limit, status=status)
        
        if not users:
            return "User not found"
        
        result = f"Found {len(users)} users:\n\n"
        result += f"{'ID':<5} {'Login name':<15} {'Name':<20} {'Email':<25} {'state':<8}\n"
        result += f"{'-'*5} {'-'*15} {'-'*20} {'-'*25} {'-'*8}\n"
        
        for user in users:
            full_name = f"{user.firstname} {user.lastname}".strip()
            if not full_name:
                full_name = user.login
            status_text = "enable" if user.status == 1 else "deactivate"
            email = user.mail[:22] + "..." if len(user.mail) > 25 else user.mail
            result += f"{user.id:<5} {user.login:<15} {full_name:<20} {email:<25} {status_text:<8}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to get user list: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_user(user_id: int) -> str:
    """
    Get details about a specific user
    
    Args:
        user_id: user ID
        
    Returns:
        User details, presented in an easy-to-read format
    """
    try:
        client = get_client()
        user_data = client.get_user(user_id)
        
        # Format user information
        result = f"User #{user_id}: {user_data.get('firstname', '')} {user_data.get('lastname', '')}\n\n"
        result += "Basic information:\n"
        result += f"- Login name: {user_data.get('login', 'N/A')}\n"
        result += f"- Email: {user_data.get('mail', 'N/A')}\n"
        result += f"- Status: {'Active' if user_data.get('status', 1) == 1 else 'Inactive'}\n"
        result += f"- Creation time: {user_data.get('created_on', 'N/A')}\n"
        
        if user_data.get('last_login_on'):
            result += f"- Last login: {user_data.get('last_login_on')}\n"
        
        # Group information
        if user_data.get('groups'):
            result += "\nGroup:\n"
            for group in user_data['groups']:
                result += f"- {group.get('name', 'N/A')}\n"
        
        # Custom fields
        if user_data.get('custom_fields'):
            result += "\nCustom field:\n"
            for field in user_data['custom_fields']:
                if field.get('value'):
                    result += f"- {field.get('name', 'N/A')}: {field.get('value', 'N/A')}\n"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to obtain user information: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def add_watcher(issue_id: int, user_id: int = None, user_name: str = None, user_login: str = None) -> str:
    """
    Add topic observer

    Args:
        issue_id: Issue ID
        user_id: user ID (choose one from user_name/user_login)
        user_name: user name (choose one from user_id/user_login)
        user_login: user login name (choose one from user_id/user_name)

    Returns:
        Operation result message
    """
    try:
        client = get_client()

        # Parse user ID
        final_user_id = user_id
        if user_name:
            final_user_id = client.find_user_id_by_name(user_name)
            if not final_user_id:
                users = client.get_available_users()
                return f'Username not found: "{user_name}"\n\nAvailable users: \n' + "\n".join([f"- {name}" for name in users['by_name'].keys()])
        elif user_login:
            final_user_id = client.find_user_id_by_login(user_login)
            if not final_user_id:
                users = client.get_available_users()
                return f'User login name not found: "{user_login}"\n\nAvailable users: \n' + "\n".join([f"- {login}" for login in users['by_login'].keys()])

        if not final_user_id:
            return "Error: Please provide user_id, user_name or user_login"

        client.add_watcher(issue_id, final_user_id)
        return f"User (ID: {final_user_id}) Add to issue #{issue_id} observer"

    except RedmineAPIError as e:
        return f"Failed to add observer: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def remove_watcher(issue_id: int, user_id: int = None, user_name: str = None, user_login: str = None) -> str:
    """
    Remove issue observer

    Args:
        issue_id: Issue ID
        user_id: user ID (choose one from user_name/user_login)
        user_name: user name (choose one from user_id/user_login)
        user_login: user login name (choose one from user_id/user_name)

    Returns:
        Operation result message
    """
    try:
        client = get_client()

        # Parse user ID
        final_user_id = user_id
        if user_name:
            final_user_id = client.find_user_id_by_name(user_name)
            if not final_user_id:
                users = client.get_available_users()
                return f'Username not found: "{user_name}"\n\nAvailable users: \n' + "\n".join([f"- {name}" for name in users['by_name'].keys()])
        elif user_login:
            final_user_id = client.find_user_id_by_login(user_login)
            if not final_user_id:
                users = client.get_available_users()
                return f'User login name not found: "{user_login}"\n\nAvailable users: \n' + "\n".join([f"- {login}" for login in users['by_login'].keys()])

        if not final_user_id:
            return "Error: Please provide user_id, user_name or user_login"

        client.remove_watcher(issue_id, final_user_id)
        return f"User (ID: {final_user_id}) from issue #{issue_id} removed from observers"

    except RedmineAPIError as e:
        return f"Failed to remove observer: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def refresh_cache() -> str:
    """
    Manually refresh enumeration values and user cache
    
    Returns:
        Refresh result message
    """
    try:
        client = get_client()
        client.refresh_cache()
        
        # Get cache information
        cache = client._load_enum_cache()
        domain = cache.get('domain', 'N/A')
        cache_time = cache.get('cache_time', 0)
        
        if cache_time > 0:
            cache_datetime = datetime.fromtimestamp(cache_time).strftime('%Y-%m-%d %H:%M:%S')
        else:
            cache_datetime = 'N/A'
        
        result = f"""Cache refreshed successfully!

Domain: {domain}
Cache time: {cache_datetime}

Cache content statistics:
- Priority: {len(cache.get('priorities', {}))} a
- Status: {len(cache.get('statuses', {}))} a
- Tracker: {len(cache.get('trackers', {}))} a
- User (name): {len(cache.get('users_by_name', {}))} a
- User (login name): {len(cache.get('users_by_login', {}))} a

Cache location: {client._cache_file}"""
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to refresh cache: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def list_issue_journals(issue_id: int, include_property_changes: bool = False) -> str:
    """
    List all notes/log records for the issue
    
    Args:
        issue_id: Issue ID
        include_property_changes: Whether to include property change records (default is no, only records with remarks will be displayed)
    
    Returns:
        List of topic notes, including Journal ID, author, time, and content
    """
    try:
        client = get_client()
        journals = client.get_issue_journals(issue_id)
        
        if not journals:
            return f"Issue #{issue_id} No notes recorded"
        
        # Filter based on parameters
        if not include_property_changes:
            # Only show records with comments
            filtered_journals = [j for j in journals if j.get('notes', '').strip()]
        else:
            filtered_journals = journals
        
        if not filtered_journals:
            return f"Issue #{issue_id} There are no qualified remark records (total {len(journals)} Pen attribute change record)"
        
        result = f"Issue #{issue_id} Note records (total {len(filtered_journals)} Pen):\n"
        result += "=" * 50 + "\n\n"
        
        for journal in filtered_journals:
            journal_id = journal.get('id', 'N/A')
            author = journal.get('user', {}).get('name', 'N/A')
            created_on = journal.get('created_on', 'N/A')
            notes = journal.get('notes', '').strip()
            private_notes = journal.get('private_notes', False)
            details = journal.get('details', [])
            
            result += f"📝 Journal #{journal_id}\n"
            result += f"   author: {author}\n"
            result += f"   time: {created_on}\n"
            if private_notes:
                result += f"   🔒 Private Note\n"
            
            if notes:
                result += f"   Remarks:\n"
                # Indent content of remarks
                for line in notes.split('\n'):
                    result += f"      {line}\n"
            
            if include_property_changes and details:
                result += f"   Property changes ({len(details)} Item):\n"
                for detail in details:
                    prop_name = detail.get('name', 'N/A')
                    old_value = detail.get('old_value', '(null)')
                    new_value = detail.get('new_value', '(null)')
                    result += f"      - {prop_name}: {old_value} → {new_value}\n"
            
            result += "\n"
        
        return result.strip()
        
    except RedmineAPIError as e:
        return f"Failed to obtain issue notes: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


@mcp.tool()
def get_journal(issue_id: int, journal_id: int) -> str:
    """
    Get details about specific comments in an issue
    
    Args:
        issue_id: Issue ID
        journal_id: Note ID (Journal ID)
    
    Returns:
        Detailed information of the note, including author, time, content, attribute changes, etc.
    """
    try:
        client = get_client()
        journals = client.get_issue_journals(issue_id)
        
        # Find the specified journal
        target_journal = None
        for journal in journals:
            if journal.get('id') == journal_id:
                target_journal = journal
                break
        
        if not target_journal:
            return f"Note not found: Issue #{issue_id} There is no Journal # in{journal_id}"
        
        journal_id = target_journal.get('id', 'N/A')
        author = target_journal.get('user', {}).get('name', 'N/A')
        author_id = target_journal.get('user', {}).get('id', 'N/A')
        created_on = target_journal.get('created_on', 'N/A')
        notes = target_journal.get('notes', '').strip()
        private_notes = target_journal.get('private_notes', False)
        details = target_journal.get('details', [])
        
        result = f"📝 Journal #{journal_id} Details\n"
        result += "=" * 50 + "\n\n"
        result += f"Issue: #{issue_id}\n"
        result += f"author: {author} (ID: {author_id})\n"
        result += f"Creation time: {created_on}\n"
        if private_notes:
            result += f"🔒 This is a private note\n"
        
        result += "\n--- Remarks ---\n"
        if notes:
            result += notes + "\n"
        else:
            result += "(No text remarks)\n"
        
        if details:
            result += f"\n---Attribute changes ({len(details)} item) ---\n"
            for detail in details:
                prop_name = detail.get('name', 'N/A')
                property_type = detail.get('property', 'N/A')
                old_value = detail.get('old_value', '(null)')
                new_value = detail.get('new_value', '(null)')
                result += f"• {prop_name} ({property_type})\n"
                result += f"  Old value: {old_value}\n"
                result += f"  New value: {new_value}\n"
        
        return result.strip()
        
    except RedmineAPIError as e:
        return f"Failed to get notes: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


# Attachment processing related constants
MAX_ATTACHMENT_SIZE_BYTES = 10 * 1024 * 1024  # 10MB
MAX_TEXT_OUTPUT_LENGTH = 50000  # Maximum number of characters for text output
DEFAULT_THUMBNAIL_SIZE = 800  # Default maximum side length of thumbnails
SUPPORTED_IMAGE_TYPES = {'image/png', 'image/jpeg', 'image/gif', 'image/webp'}

# Explicitly supported text MIME types
TEXT_MIME_TYPES = {
    'text/plain', 'text/html', 'text/css', 'text/csv',
    'text/xml', 'text/javascript', 'text/markdown',
    'application/json', 'application/xml',
    'application/javascript', 'application/x-yaml',
    'application/x-sh', 'application/sql',
}

# The file type corresponding to the file extension (used for backup judgment when content_type is inaccurate)
TEXT_EXTENSIONS = {
    '.txt', '.csv', '.json', '.xml', '.yaml', '.yml',
    '.html', '.htm', '.css', '.js', '.ts', '.md',
    '.sh', '.bash', '.zsh', '.py', '.rb', '.php',
    '.java', '.go', '.rs', '.sql', '.log', '.ini',
    '.cfg', '.conf', '.toml', '.env', '.gitignore',
}

OFFICE_EXTENSIONS = {
    '.pdf': 'pdf',
    '.docx': 'docx',
    '.xlsx': 'xlsx',
    '.pptx': 'pptx',
}

# Old Office formats (direct parsing is not supported)
LEGACY_OFFICE_EXTENSIONS = {'.doc', '.xls', '.ppt'}


@mcp.tool()
def get_attachment_image(
    attachment_id: int,
    thumbnail: bool = True,
    max_size: int = DEFAULT_THUMBNAIL_SIZE
):
    """
    Download the Redmine attachment image for AI visual analysis
    
    Args:
        attachment_id: attachment ID
        thumbnail: whether to generate thumbnails to reduce token consumption (default True)
        max_size: Maximum side length of thumbnail (default 800 pixels, only effective when thumbnail=True)
    
    Returns:
        Image content (for AI visual interpretation) or error message
    """
    try:
        from io import BytesIO
        from PIL import Image as PILImage
        
        client = get_client()
        
        # Download attachment
        image_data, attachment_info = client.download_attachment(attachment_id)
        
        filename = attachment_info.get('filename', 'unknown')
        content_type = attachment_info.get('content_type', '')
        filesize = attachment_info.get('filesize', 0)
        
        # Check if it is an image type
        if content_type not in SUPPORTED_IMAGE_TYPES:
            return f"appendix #{attachment_id} ({filename}) is not a supported image format\n" \
                   f"type: {content_type}\n" \
                   f"Supported formats: {', '.join(SUPPORTED_IMAGE_TYPES)}"
        
        # Check file size
        if filesize > MAX_ATTACHMENT_SIZE_BYTES:
            return f"appendix #{attachment_id} ({filename}) file too large\n" \
                   f"size: {filesize / 1024 / 1024:.2f} MB\n" \
                   f"limit: {MAX_ATTACHMENT_SIZE_BYTES / 1024 / 1024:.0f} MB"
        
        # Process pictures
        img = PILImage.open(BytesIO(image_data))
        original_size = img.size
        
        # Conversion modes (handling RGBA, P, etc. modes)
        if img.mode in ('RGBA', 'P'):
            # Create a white background
            background = PILImage.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Thumbnail processing
        if thumbnail and (img.width > max_size or img.height > max_size):
            img.thumbnail((max_size, max_size), PILImage.Resampling.LANCZOS)
            resized = True
        else:
            resized = False
        
        # Output as JPEG (smaller file size)
        output_buffer = BytesIO()
        img.save(output_buffer, format='JPEG', quality=85, optimize=True)
        output_data = output_buffer.getvalue()
        
        # Logging processing information (via logging)
        import logging
        logger = logging.getLogger(__name__)
        if resized:
            logger.info(
                f"picture #{attachment_id} ({filename}): "
                f"{original_size[0]}x{original_size[1]} → {img.size[0]}x{img.size[1]}, "
                f"{len(image_data)} → {len(output_data)} bytes"
            )
        
        return Image(data=output_data, format="jpeg")
        
    except RedmineAPIError as e:
        return f"Failed to obtain attached image: {str(e)}"
    except Exception as e:
        return f"Handling image errors: {str(e)}"


def _get_file_extension(filename: str) -> str:
    """Get the file extension (lowercase)"""
    return os.path.splitext(filename)[1].lower()


def _extract_pdf_text(data: bytes) -> str:
    """Extract text from PDF"""
    from io import BytesIO
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"--- No. {i + 1} Page ---\n{text.strip()}")

    if not pages:
        return "(PDF cannot extract text content, it may be a scanned image PDF)"
    return "\n\n".join(pages)


def _extract_docx_text(data: bytes) -> str:
    """Extract text from Word (.docx)"""
    from io import BytesIO
    from docx import Document

    doc = Document(BytesIO(data))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # Also extract table content
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            paragraphs.append("\n[Form]\n" + "\n".join(rows))

    if not paragraphs:
        return "(Word file has no text content)"
    return "\n\n".join(paragraphs)


def _extract_xlsx_text(data: bytes) -> str:
    """Extract table text from Excel (.xlsx)"""
    from io import BytesIO
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(c for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"--- Worksheet: {sheet_name} ---\n" + "\n".join(rows))

    wb.close()
    if not sheets:
        return "(Excel file has no content)"
    return "\n\n".join(sheets)


def _extract_pptx_text(data: bytes) -> str:
    """Extract text from PowerPoint (.pptx)"""
    from io import BytesIO
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    slides = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            slides.append(f"--- Slides {i + 1} ---\n" + "\n".join(texts))

    if not slides:
        return "(Presentation has no text content)"
    return "\n\n".join(slides)


def _try_decode_text(data: bytes) -> str | None:
    """Attempts to decode bytes into text, returning None on failure"""
    for encoding in ('utf-8', 'utf-8-sig', 'big5', 'gb2312', 'shift_jis', 'latin-1'):
        try:
            text = data.decode(encoding)
            # Check for too many non-printable characters (binary feature)
            non_printable = sum(1 for c in text[:1000] if not c.isprintable() and c not in '\n\r\t')
            if non_printable / max(len(text[:1000]), 1) > 0.1:
                continue
            return text
        except (UnicodeDecodeError, ValueError):
            continue
    return None


@mcp.tool()
def get_attachment_text(
    attachment_id: int,
    max_length: int = MAX_TEXT_OUTPUT_LENGTH
) -> str:
    """
    Read the text content of Redmine attachments

    Supported formats:
    - File types: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx)
    - Text: TXT, JSON, XML, CSV, HTML, Markdown, program code, etc.
    - Others: Automatically try to read in plain text, binary files will be prompted not supported

    Note: Please use get_attachment_image() for images

    Args:
        attachment_id: attachment ID
        max_length: The maximum number of characters for text output (default 50000)

    Returns:
        Attachment text or error message
    """
    try:
        client = get_client()

        # Download attachment
        file_data, attachment_info = client.download_attachment(attachment_id)

        filename = attachment_info.get('filename', 'unknown')
        content_type = attachment_info.get('content_type', '')
        filesize = attachment_info.get('filesize', 0)
        ext = _get_file_extension(filename)

        # Check file size
        if filesize > MAX_ATTACHMENT_SIZE_BYTES:
            return f"appendix #{attachment_id} ({filename}) file too large\n" \
                   f"size: {filesize / 1024 / 1024:.2f} MB\n" \
                   f"limit: {MAX_ATTACHMENT_SIZE_BYTES / 1024 / 1024:.0f} MB"

        # Image type → prompt to use get_attachment_image
        if content_type in SUPPORTED_IMAGE_TYPES:
            return f"appendix #{attachment_id} ({filename}) is the image file\n" \
                   f"Please use get_attachment_image({attachment_id}) for visual analysis"

        # Old Office formatting tips
        if ext in LEGACY_OFFICE_EXTENSIONS:
            new_ext = ext + 'x'
            return f"appendix #{attachment_id} ({filename}) is an older Office format ({ext})\n" \
                   f"Currently only new formats are supported ({new_ext}), please convert the file to the new format and upload it again."

        # Office file types
        office_type = OFFICE_EXTENSIONS.get(ext)
        extracted_text = None
        format_label = ""

        try:
            if office_type == 'pdf' or content_type == 'application/pdf':
                extracted_text = _extract_pdf_text(file_data)
                format_label = "PDF"
            elif office_type == 'docx' or content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                extracted_text = _extract_docx_text(file_data)
                format_label = "Word"
            elif office_type == 'xlsx' or content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                extracted_text = _extract_xlsx_text(file_data)
                format_label = "Excel"
            elif office_type == 'pptx' or content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                extracted_text = _extract_pptx_text(file_data)
                format_label = "PowerPoint"
        except Exception as e:
            return f"appendix #{attachment_id} ({filename}) parsing failed\n" \
                   f"type: {content_type}\n" \
                   f"mistake: {str(e)}"

        if extracted_text is not None:
            pass  # It has been parsed by Office above
        elif content_type in TEXT_MIME_TYPES or content_type.startswith('text/') or ext in TEXT_EXTENSIONS:
            extracted_text = _try_decode_text(file_data)
            format_label = "text"
            if extracted_text is None:
                return f"appendix #{attachment_id} ({filename}) cannot be decoded as literal \n" \
                       f"type: {content_type}"
        else:
            # Unknown type → try plain text read
            extracted_text = _try_decode_text(file_data)
            if extracted_text is not None:
                format_label = "Text (automatically detected)"
            else:
                return f"appendix #{attachment_id} ({filename}) is a binary file and cannot read the text content\n" \
                       f"type: {content_type}\n" \
                       f"size: {filesize / 1024:.2f} KB"

        # Truncate overly long content
        truncated = False
        if len(extracted_text) > max_length:
            extracted_text = extracted_text[:max_length]
            truncated = True

        header = f"appendix #{attachment_id} ({filename}) - {format_label} content"
        header += f"\n{'=' * 50}\n"
        result = header + extracted_text

        if truncated:
            result += f"\n\n... (Content has been truncated, original length exceeds {max_length} characters)"

        return result

    except RedmineAPIError as e:
        return f"Failed to obtain attachment content: {str(e)}"
    except Exception as e:
        return f"Handling attachment errors: {str(e)}"


@mcp.tool()
def get_attachment_info(attachment_id: int) -> str:
    """
    Get attachment details (without downloading file contents)
    
    Args:
        attachment_id: attachment ID
    
    Returns:
        Attachment details
    """
    try:
        client = get_client()
        attachment = client.get_attachment(attachment_id)
        
        filesize = attachment.get('filesize', 0)
        size_text = f"{filesize / 1024 / 1024:.2f} MB" if filesize >= 1024 * 1024 else f"{filesize / 1024:.2f} KB"
        
        result = f"""appendix #{attachment.get('id')} Details
==================================================

File name: {attachment.get('filename', 'N/A')}
File size: {size_text}
File type: {attachment.get('content_type', 'N/A')}
illustrate: {attachment.get('description', '(no explanation)')}

Uploader: {attachment.get('author', {}).get('name', 'N/A')}
Upload time: {attachment.get('created_on', 'N/A')}

Download link: {attachment.get('content_url', 'N/A')}"""
        
        # Provide action suggestions based on file type
        content_type = attachment.get('content_type', '')
        att_id = attachment.get('id')
        att_filename = attachment.get('filename', '')
        att_ext = _get_file_extension(att_filename)

        if content_type in SUPPORTED_IMAGE_TYPES:
            result += f"\n\n💡 This is an image file, you can use get_attachment_image({att_id}) for visual analysis"
        elif (content_type == 'application/pdf' or att_ext == '.pdf'
              or content_type in TEXT_MIME_TYPES or content_type.startswith('text/')
              or att_ext in TEXT_EXTENSIONS or att_ext in OFFICE_EXTENSIONS):
            result += f"\n\n💡 can use get_attachment_text({att_id}) Read text content"
        else:
            result += f"\n\n💡 You can try get_attachment_text({att_id}) Read content (unknown formats will automatically try plain text reading)"
        
        return result
        
    except RedmineAPIError as e:
        return f"Failed to obtain attachment information: {str(e)}"
    except Exception as e:
        return f"System error: {str(e)}"


def main():
    """MCP server main entry point"""
    parser = argparse.ArgumentParser(
        description='Redmine MCP server',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
Example:
  redmine-mcp #Default stdio mode
  redmine-mcp --transport sse # SSE mode (default port 8000)
  redmine-mcp --transport sse --port 3000 --host 127.0.0.1
        '''
    )
    parser.add_argument(
        '--transport', '-t',
        choices=['stdio', 'sse'],
        default=None,
        help='Transmission mode: stdio (default) or sse'
    )
    parser.add_argument(
        '--host', '-H',
        default=None,
        help='SSE mode binding address (default: 0.0.0.0)'
    )
    parser.add_argument(
        '--port', '-p',
        type=int,
        default=None,
        help='SSE mode listening port (default: 8000)'
    )

    args = parser.parse_args()
    config = get_config()

    # Command line parameters take precedence over environment variables
    transport = args.transport or config.transport

    if transport == 'sse':
        host = args.host or config.sse_host
        port = args.port or config.sse_port
        # Set SSE server parameters
        mcp.settings.host = host
        mcp.settings.port = port
        print(f"Start SSE mode: http://{host}:{port}", flush=True)
        mcp.run('sse')
    else:
        mcp.run('stdio')


if __name__ == "__main__":
    main()